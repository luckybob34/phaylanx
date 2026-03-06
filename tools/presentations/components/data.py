"""
Data component renderers — stat-grid, card-grid, data-table, highlight-grid.
"""

from __future__ import annotations

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .helpers import (
    add_textbox, add_rect, resolve_color, apply_font,
    distribute_x, center_x,
)
from .typography import render_slide_chrome


def render_stat_grid(slide, slide_data: dict, theme: dict):
    """Render a stat grid — 3-6 metric cards in columns."""
    y_cursor = render_slide_chrome(slide, slide_data, theme)

    stats = slide_data.get("stats", [])
    if not stats:
        return

    fonts = theme.get("fonts", {})
    stat_font = fonts.get("stat", {})
    body_font = fonts.get("body", {})
    comp = theme.get("components", {}).get("stat_grid", {})
    margins = theme.get("margins", {})
    slide_w = theme.get("slide", {}).get("width_in", 13.333)

    count = len(stats)
    max_cols = min(comp.get("max_columns", 3), count)
    card_width = 2.8 if count <= 3 else 2.2
    gap = 0.3
    positions = distribute_x(count, card_width, gap, slide_width=slide_w)

    card_height = 2.0
    card_top = y_cursor + 0.3

    for i, stat in enumerate(stats):
        x = positions[i]

        # Card background
        card = add_rect(slide, x, card_top, card_width, card_height,
                        "background_warm", theme, corner_radius_emu=91440)

        # Accent top stripe
        stripe = add_rect(slide, x + 0.1, card_top + 0.1, card_width - 0.2, 0.06,
                          "accent", theme)

        # Value
        tb = add_textbox(slide, x + 0.2, card_top + 0.35, card_width - 0.4, 0.8)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = stat.get("value", "")
        apply_font(run, stat_font, theme, "size_medium")

        # Label
        tb = add_textbox(slide, x + 0.2, card_top + 1.25, card_width - 0.4, 0.6)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = stat.get("label", "")
        run.font.name = body_font.get("name", "Lato")
        run.font.size = Pt(comp.get("label_size", 11))
        run.font.color.rgb = resolve_color("primary_light", theme)


def render_card_grid(slide, slide_data: dict, theme: dict):
    """Render a card grid — auto-fit content cards."""
    y_cursor = render_slide_chrome(slide, slide_data, theme)

    cards = slide_data.get("cards", [])
    if not cards:
        return

    fonts = theme.get("fonts", {})
    heading_font = fonts.get("heading", {})
    body_font = fonts.get("body", {})
    comp = theme.get("components", {}).get("card_grid", {})
    margins = theme.get("margins", {})
    slide_w = theme.get("slide", {}).get("width_in", 13.333)

    count = len(cards)
    columns = slide_data.get("columns", min(comp.get("max_columns", 4), count))
    rows = (count + columns - 1) // columns

    gap = comp.get("card_gap_in", 0.2)
    total_width = slide_w - margins.get("left", 0.66) - margins.get("right", 0.66)
    card_width = (total_width - gap * (columns - 1)) / columns
    card_height = 2.2 if rows <= 2 else 1.8

    start_x = margins.get("left", 0.66)
    card_top = y_cursor + 0.3

    for i, card in enumerate(cards):
        col = i % columns
        row = i // columns
        x = start_x + col * (card_width + gap)
        y = card_top + row * (card_height + gap)

        # Card background
        bg = add_rect(slide, x, y, card_width, card_height,
                      "background_warm", theme,
                      corner_radius_emu=comp.get("card_corner_radius_emu", 91440))

        # Icon (emoji)
        icon = card.get("icon", "")
        text_y = y + 0.15
        if icon:
            tb = add_textbox(slide, x + 0.15, text_y, 0.4, 0.4)
            tf = tb.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = icon
            run.font.size = Pt(18)
            text_y += 0.4

        # Card title
        tb = add_textbox(slide, x + 0.15, text_y, card_width - 0.3, 0.4)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = card.get("title", "")
        run.font.name = heading_font.get("name", "Source Serif Pro SemiBold")
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = resolve_color("primary", theme)

        # Card body
        body = card.get("body", "")
        if body:
            tb = add_textbox(slide, x + 0.15, text_y + 0.45, card_width - 0.3, card_height - (text_y - y) - 0.6)
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = body
            run.font.name = body_font.get("name", "Lato")
            run.font.size = Pt(11)
            run.font.color.rgb = resolve_color("primary_light", theme)


def render_data_table(slide, slide_data: dict, theme: dict):
    """Render a data table with styled header row."""
    y_cursor = render_slide_chrome(slide, slide_data, theme)

    headers = slide_data.get("headers", [])
    rows = slide_data.get("rows", [])
    if not headers:
        return

    fonts = theme.get("fonts", {})
    body_font = fonts.get("body", {})
    comp = theme.get("components", {}).get("data_table", {})
    margins = theme.get("margins", {})
    slide_w = theme.get("slide", {}).get("width_in", 13.333)

    left = margins.get("left", 0.66)
    total_width = slide_w - left - margins.get("right", 0.66)
    col_count = len(headers)
    col_width = total_width / col_count
    row_height = 0.45
    header_height = 0.45

    table_top = y_cursor + 0.3

    # Add table shape
    table_shape = slide.shapes.add_table(
        len(rows) + 1, col_count,
        Inches(left), Inches(table_top),
        Inches(total_width), Inches(header_height + row_height * len(rows))
    )
    table = table_shape.table

    # Style header
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.name = body_font.get("name", "Lato")
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = resolve_color(
                    comp.get("header_text", "text_on_dark"), theme
                )
        cell.fill.solid()
        cell.fill.fore_color.rgb = resolve_color(
            comp.get("header_bg", "accent"), theme
        )

    # Style data rows
    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = body_font.get("name", "Lato")
                    run.font.size = Pt(11)
                    run.font.color.rgb = resolve_color("primary", theme)

            # Alternate row background
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = resolve_color(
                    comp.get("alt_row_bg", "background_warm"), theme
                )


def render_highlight_grid(slide, slide_data: dict, theme: dict):
    """Render a highlight/feature grid — same as card grid but with accent stripe."""
    # Reuse card_grid but with accent top bar
    y_cursor = render_slide_chrome(slide, slide_data, theme)

    items = slide_data.get("cards", [])
    if not items:
        return

    fonts = theme.get("fonts", {})
    heading_font = fonts.get("heading", {})
    body_font = fonts.get("body", {})
    margins = theme.get("margins", {})
    slide_w = theme.get("slide", {}).get("width_in", 13.333)

    count = len(items)
    columns = slide_data.get("columns", min(4, count))
    gap = 0.25
    total_width = slide_w - margins.get("left", 0.66) - margins.get("right", 0.66)
    card_width = (total_width - gap * (columns - 1)) / columns
    card_height = 2.0
    start_x = margins.get("left", 0.66)
    card_top = y_cursor + 0.3

    # Palette cycling for accent stripes
    palette_keys = ["accent", "slate_blue", "sage", "gold", "light_blue"]

    for i, item in enumerate(items):
        col = i % columns
        row = i // columns
        x = start_x + col * (card_width + gap)
        y = card_top + row * (card_height + gap)
        color_key = palette_keys[i % len(palette_keys)]

        # Card background
        add_rect(slide, x, y, card_width, card_height,
                 "background_warm", theme, corner_radius_emu=91440)

        # Accent top stripe
        add_rect(slide, x, y, card_width, 0.06, color_key, theme)

        # Icon
        icon = item.get("icon", "")
        text_y = y + 0.2
        if icon:
            tb = add_textbox(slide, x + 0.15, text_y, 0.4, 0.4)
            p = tb.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = icon
            run.font.size = Pt(18)
            text_y += 0.4

        # Title
        tb = add_textbox(slide, x + 0.15, text_y, card_width - 0.3, 0.35)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = item.get("title", "")
        run.font.name = heading_font.get("name", "Source Serif Pro SemiBold")
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = resolve_color("primary", theme)

        # Description
        desc = item.get("body", item.get("desc", ""))
        if desc:
            tb = add_textbox(slide, x + 0.15, text_y + 0.4, card_width - 0.3, card_height - 1.0)
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = desc
            run.font.name = body_font.get("name", "Lato")
            run.font.size = Pt(10)
            run.font.color.rgb = resolve_color("primary_light", theme)
