"""
Structure & process renderers — layer-stack, hub-spoke, process-loop, comparison.
"""

from __future__ import annotations

import math

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .helpers import (
    add_textbox, add_rect, add_oval, add_line, add_arrow_right,
    resolve_color, apply_font, center_x, point_on_circle,
)
from .typography import render_slide_chrome


def render_layer_stack(slide, slide_data: dict, theme: dict):
    """Render tiered architecture rows."""
    y_cursor = render_slide_chrome(slide, slide_data, theme)

    layers = slide_data.get("layers", [])
    if not layers:
        return

    fonts = theme.get("fonts", {})
    body_font = fonts.get("body", {})
    heading_font = fonts.get("heading", {})
    comp = theme.get("components", {}).get("layer_stack", {})
    margins = theme.get("margins", {})
    slide_w = theme.get("slide", {}).get("width_in", 13.333)

    row_h = comp.get("row_height_in", 0.9)
    row_gap = comp.get("row_gap_in", 0.12)
    layer_colors = comp.get("layer_colors", ["accent", "slate_blue", "sage", "gold", "light_blue"])

    left = margins.get("left", 0.66)
    total_width = slide_w - 2 * left
    label_width = 2.2
    items_left = left + label_width + 0.2
    items_width = total_width - label_width - 0.2

    row_top = y_cursor + 0.4

    for i, layer in enumerate(layers):
        y = row_top + i * (row_h + row_gap)
        color = layer_colors[i % len(layer_colors)]

        # Label band on left
        add_rect(slide, left, y, label_width, row_h, color, theme,
                 corner_radius_emu=45720)
        tb = add_textbox(slide, left + 0.1, y, label_width - 0.2, row_h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        # Vertical centre
        p.space_before = Pt((row_h * 72 - 12) / 2)
        run = p.add_run()
        run.text = layer.get("label", "")
        run.font.name = body_font.get("name", "Lato")
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = resolve_color("text_on_dark", theme)

        # Item chips
        items = layer.get("items", [])
        if items:
            item_count = len(items)
            chip_gap = 0.15
            chip_w = (items_width - chip_gap * (item_count - 1)) / item_count

            for j, item in enumerate(items):
                chip_x = items_left + j * (chip_w + chip_gap)
                add_rect(slide, chip_x, y, chip_w, row_h,
                         "background_warm", theme, corner_radius_emu=45720)
                tb = add_textbox(slide, chip_x + 0.1, y, chip_w - 0.2, row_h)
                tf = tb.text_frame
                tf.word_wrap = True
                tf.auto_size = None
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.space_before = Pt((row_h * 72 - 11) / 2)
                run = p.add_run()
                run.text = item
                run.font.name = body_font.get("name", "Lato")
                run.font.size = Pt(10)
                run.font.color.rgb = resolve_color("primary", theme)


def render_hub_spoke(slide, slide_data: dict, theme: dict):
    """Render a hub-spoke diagram — centre node with radial satellites."""
    y_cursor = render_slide_chrome(slide, slide_data, theme)

    spokes = slide_data.get("spokes", [])
    center_label = slide_data.get("center", "")
    if not spokes:
        return

    fonts = theme.get("fonts", {})
    body_font = fonts.get("body", {})
    comp = theme.get("components", {}).get("hub_spoke", {})
    slide_w = theme.get("slide", {}).get("width_in", 13.333)
    bottom = theme.get("margins", {}).get("footer_bottom", 7.0)

    center_r = comp.get("center_radius_in", 0.8)
    spoke_r = comp.get("spoke_radius_in", 2.5)
    node_r = comp.get("node_radius_in", 0.5)

    # Centre of the diagram
    available_h = bottom - y_cursor - 0.5
    cx = slide_w / 2
    cy = y_cursor + available_h / 2

    # Centre node
    add_oval(slide, cx - center_r, cy - center_r, center_r * 2, center_r * 2,
             "accent", theme)
    tb = add_textbox(slide, cx - center_r + 0.1, cy - 0.2, (center_r - 0.1) * 2, 0.5)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = center_label
    run.font.name = body_font.get("name", "Lato")
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = resolve_color("text_on_dark", theme)

    # Spoke nodes
    count = len(spokes)
    palette = ["slate_blue", "sage", "gold", "light_blue", "accent_light",
               "slate_blue", "sage", "gold"]

    for i, spoke_text in enumerate(spokes):
        angle = (360 / count) * i
        sx, sy = point_on_circle(cx, cy, spoke_r, angle)

        # Connecting line
        add_line(slide, cx, cy, sx, sy, "primary_light", theme, width_pt=1.0)

        # Spoke node
        color = palette[i % len(palette)]
        add_oval(slide, sx - node_r, sy - node_r, node_r * 2, node_r * 2,
                 color, theme)
        tb = add_textbox(slide, sx - node_r + 0.05, sy - 0.15, (node_r - 0.05) * 2, 0.35)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = spoke_text
        run.font.name = body_font.get("name", "Lato")
        run.font.size = Pt(8)
        run.font.bold = True
        run.font.color.rgb = resolve_color("text_on_dark", theme)


def render_process_loop(slide, slide_data: dict, theme: dict):
    """Render a circular process loop with nodes around a centre."""
    y_cursor = render_slide_chrome(slide, slide_data, theme)

    nodes = slide_data.get("nodes", [])
    center_label = slide_data.get("center_label", "")
    if not nodes:
        return

    fonts = theme.get("fonts", {})
    body_font = fonts.get("body", {})
    comp = theme.get("components", {}).get("process_loop", {})
    slide_w = theme.get("slide", {}).get("width_in", 13.333)
    bottom = theme.get("margins", {}).get("footer_bottom", 7.0)

    ring_r = comp.get("ring_radius_in", 2.0)
    node_r = comp.get("node_radius_in", 0.45)
    center_r = comp.get("center_radius_in", 0.6)

    available_h = bottom - y_cursor - 0.5
    cx = slide_w / 2
    cy = y_cursor + available_h / 2

    # Centre circle
    add_oval(slide, cx - center_r, cy - center_r, center_r * 2, center_r * 2,
             "accent", theme)
    if center_label:
        tb = add_textbox(slide, cx - center_r + 0.05, cy - 0.2, (center_r - 0.05) * 2, 0.45)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = center_label
        run.font.name = body_font.get("name", "Lato")
        run.font.size = Pt(8)
        run.font.bold = True
        run.font.color.rgb = resolve_color("text_on_dark", theme)

    # Ring nodes
    count = len(nodes)
    palette = ["accent", "slate_blue", "sage", "gold", "light_blue", "accent_dark"]

    for i, node_text in enumerate(nodes):
        angle = (360 / count) * i
        nx, ny = point_on_circle(cx, cy, ring_r, angle)
        color = palette[i % len(palette)]

        # Connecting line from previous node to this one
        if i > 0:
            prev_angle = (360 / count) * (i - 1)
            px, py = point_on_circle(cx, cy, ring_r, prev_angle)
            add_line(slide, px, py, nx, ny, "primary_light", theme, width_pt=1.5)
        # Close the loop
        if i == count - 1:
            first_x, first_y = point_on_circle(cx, cy, ring_r, 0)
            add_line(slide, nx, ny, first_x, first_y, "primary_light", theme, width_pt=1.5)

        # Node circle
        add_oval(slide, nx - node_r, ny - node_r, node_r * 2, node_r * 2,
                 color, theme)
        tb = add_textbox(slide, nx - node_r + 0.03, ny - 0.15,
                         (node_r - 0.03) * 2, 0.35)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = node_text
        run.font.name = body_font.get("name", "Lato")
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = resolve_color("text_on_dark", theme)


def render_comparison(slide, slide_data: dict, theme: dict):
    """Render a before/after comparison with side-by-side panels."""
    y_cursor = render_slide_chrome(slide, slide_data, theme)

    before = slide_data.get("before", {})
    after = slide_data.get("after", {})
    if not before and not after:
        return

    fonts = theme.get("fonts", {})
    body_font = fonts.get("body", {})
    heading_font = fonts.get("heading", {})
    comp = theme.get("components", {}).get("comparison", {})
    margins = theme.get("margins", {})
    slide_w = theme.get("slide", {}).get("width_in", 13.333)
    bottom = margins.get("footer_bottom", 7.0)

    panel_gap = comp.get("panel_gap_in", 0.4)
    left = margins.get("left", 0.66)
    total_width = slide_w - 2 * left
    panel_width = (total_width - panel_gap) / 2
    panel_height = bottom - y_cursor - 0.8
    panel_top = y_cursor + 0.4

    # Arrow in the middle
    arrow_size = 0.35
    arrow_x = left + panel_width + (panel_gap - arrow_size) / 2
    arrow_y = panel_top + panel_height / 2 - arrow_size / 2
    add_arrow_right(slide, arrow_x, arrow_y, arrow_size, arrow_size,
                    "accent", theme)

    # Render each panel
    for i, (panel_data, x_pos, color_key) in enumerate([
        (before, left, comp.get("before_color", "primary_light")),
        (after, left + panel_width + panel_gap, comp.get("after_color", "accent")),
    ]):
        if not panel_data:
            continue

        # Panel background
        add_rect(slide, x_pos, panel_top, panel_width, panel_height,
                 "background_warm", theme, corner_radius_emu=91440)

        # Colour accent top bar
        add_rect(slide, x_pos, panel_top, panel_width, 0.06, color_key, theme)

        # Panel title
        title = panel_data.get("title", "")
        tb = add_textbox(slide, x_pos + 0.2, panel_top + 0.25, panel_width - 0.4, 0.4)
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = heading_font.get("name", "Source Serif Pro SemiBold")
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = resolve_color(color_key, theme)

        # Items as bullet list
        items = panel_data.get("items", [])
        if items:
            tb = add_textbox(slide, x_pos + 0.2, panel_top + 0.8,
                             panel_width - 0.4, panel_height - 1.2)
            tf = tb.text_frame
            tf.word_wrap = True
            first = True
            for item in items:
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()
                p.level = 0
                p.space_before = Pt(6)
                run = p.add_run()
                run.text = f"• {item}"
                run.font.name = body_font.get("name", "Lato")
                run.font.size = Pt(11)
                run.font.color.rgb = resolve_color("primary", theme)
