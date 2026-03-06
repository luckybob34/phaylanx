"""
Flow & sequence renderers — step-flow, funnel, timeline, ascend.
"""

from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .helpers import (
    add_textbox, add_rect, add_oval, add_arrow_right,
    resolve_color, apply_font, distribute_x,
)
from .typography import render_slide_chrome


def render_step_flow(slide, slide_data: dict, theme: dict):
    """Render horizontal numbered step nodes with arrows."""
    y_cursor = render_slide_chrome(slide, slide_data, theme)

    steps = slide_data.get("steps", [])
    if not steps:
        return

    fonts = theme.get("fonts", {})
    body_font = fonts.get("body", {})
    heading_font = fonts.get("heading", {})
    comp = theme.get("components", {}).get("step_flow", {})
    slide_w = theme.get("slide", {}).get("width_in", 13.333)

    count = len(steps)
    node_w = comp.get("node_width_in", 1.8)
    node_h = comp.get("node_height_in", 1.2)
    arrow_w = comp.get("arrow_width_in", 0.4)

    # Calculate positions
    total_w = count * node_w + (count - 1) * arrow_w
    start_x = (slide_w - total_w) / 2
    node_top = y_cursor + 0.5

    for i, step in enumerate(steps):
        x = start_x + i * (node_w + arrow_w)

        # Node rectangle
        add_rect(slide, x, node_top, node_w, node_h,
                 "accent", theme, corner_radius_emu=91440)

        # Step number circle
        num_size = 0.35
        num = add_oval(slide, x + (node_w - num_size) / 2, node_top + 0.1,
                       num_size, num_size, "accent_dark", theme)
        tf = num.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(step.get("number", i + 1))
        run.font.name = body_font.get("name", "Lato")
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = resolve_color("text_on_dark", theme)

        # Step title
        tb = add_textbox(slide, x + 0.1, node_top + 0.5, node_w - 0.2, 0.3)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = step.get("title", "")
        run.font.name = body_font.get("name", "Lato")
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = resolve_color("text_on_dark", theme)

        # Step description below node
        desc = step.get("desc", "")
        if desc:
            tb = add_textbox(slide, x, node_top + node_h + 0.15, node_w, 0.8)
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = desc
            run.font.name = body_font.get("name", "Lato")
            run.font.size = Pt(9)
            run.font.color.rgb = resolve_color("primary_light", theme)

        # Arrow between nodes
        if i < count - 1:
            arrow_x = x + node_w + 0.02
            arrow_y = node_top + (node_h - 0.2) / 2
            add_arrow_right(slide, arrow_x, arrow_y, arrow_w - 0.04, 0.2,
                            "accent_light", theme)


def render_funnel(slide, slide_data: dict, theme: dict):
    """Render funnel bars — decaying horizontal bars."""
    y_cursor = render_slide_chrome(slide, slide_data, theme)

    bars = slide_data.get("bars", [])
    if not bars:
        return

    fonts = theme.get("fonts", {})
    body_font = fonts.get("body", {})
    stat_font = fonts.get("stat", {})
    comp = theme.get("components", {}).get("funnel", {})
    margins = theme.get("margins", {})
    slide_w = theme.get("slide", {}).get("width_in", 13.333)

    max_bar_w = comp.get("max_bar_width_in", 10.0)
    bar_h = comp.get("bar_height_in", 0.55)
    bar_gap = comp.get("bar_gap_in", 0.15)
    bar_colors = comp.get("bar_colors", ["accent", "gold", "sage", "slate_blue", "light_blue"])

    # Centre the funnel
    funnel_left = (slide_w - max_bar_w) / 2
    bar_top = y_cursor + 0.4

    for i, bar in enumerate(bars):
        pct = bar.get("pct", 100)
        bar_w = max_bar_w * (pct / 100)
        color_key = bar_colors[i % len(bar_colors)]

        # Centre each bar
        x = (slide_w - bar_w) / 2
        y = bar_top + i * (bar_h + bar_gap)

        # Bar rectangle
        add_rect(slide, x, y, bar_w, bar_h, color_key, theme,
                 corner_radius_emu=45720)

        # Label text inside bar
        tb = add_textbox(slide, x + 0.15, y, bar_w - 0.3, bar_h)
        tf = tb.text_frame
        tf.word_wrap = False
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        # Label
        run = p.add_run()
        run.text = bar.get("label", "")
        run.font.name = body_font.get("name", "Lato")
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = resolve_color("text_on_dark", theme)

        # Value on right side
        val = bar.get("value", "")
        if val:
            tb2 = add_textbox(slide, x + 0.15, y, bar_w - 0.3, bar_h)
            tf2 = tb2.text_frame
            tf2.word_wrap = False
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.RIGHT
            run2 = p2.add_run()
            run2.text = str(val)
            run2.font.name = body_font.get("name", "Lato")
            run2.font.size = Pt(11)
            run2.font.bold = True
            run2.font.color.rgb = resolve_color("text_on_dark", theme)

        # Centre text vertically (approximate via padding)
        for tf_item in [tf, tb2.text_frame if val else None]:
            if tf_item:
                for para in tf_item.paragraphs:
                    para.space_before = Pt(bar_h * 36 / 2 - 6)  # rough vertical centre


def render_timeline(slide, slide_data: dict, theme: dict):
    """Render a horizontal timeline with status dots."""
    y_cursor = render_slide_chrome(slide, slide_data, theme)

    events = slide_data.get("events", [])
    if not events:
        return

    fonts = theme.get("fonts", {})
    body_font = fonts.get("body", {})
    heading_font = fonts.get("heading", {})
    comp = theme.get("components", {}).get("timeline", {})
    slide_w = theme.get("slide", {}).get("width_in", 13.333)
    margins = theme.get("margins", {})

    status_colors = comp.get("status_colors", {
        "complete": "6A9E98",
        "active": "E55F4C",
        "upcoming": "9B9B9B",
    })

    count = len(events)
    left_margin = margins.get("left", 0.66) + 0.5
    right_margin = margins.get("right", 0.66) + 0.5
    line_width = slide_w - left_margin - right_margin
    line_y = y_cursor + 1.2
    dot_r = 0.12

    # Timeline line
    from .helpers import add_line
    add_line(slide, left_margin, line_y, left_margin + line_width, line_y,
             "primary_light", theme, width_pt=2)

    # Events
    spacing = line_width / (count - 1) if count > 1 else 0
    for i, event in enumerate(events):
        x = left_margin + i * spacing
        status = event.get("status", "upcoming")
        color_hex = status_colors.get(status, "9B9B9B")

        # Dot
        dot = add_oval(slide, x - dot_r, line_y - dot_r, dot_r * 2, dot_r * 2,
                       color_hex, theme)

        # Date above
        tb = add_textbox(slide, x - 0.8, line_y - 0.55, 1.6, 0.3)
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = event.get("date", "")
        run.font.name = body_font.get("name", "Lato")
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = resolve_color(color_hex, theme)

        # Title below
        tb = add_textbox(slide, x - 0.9, line_y + 0.25, 1.8, 0.35)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = event.get("title", "")
        run.font.name = heading_font.get("name", "Source Serif Pro SemiBold")
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = resolve_color("primary", theme)

        # Desc below title
        desc = event.get("desc", "")
        if desc:
            tb = add_textbox(slide, x - 0.9, line_y + 0.6, 1.8, 0.6)
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = desc
            run.font.name = body_font.get("name", "Lato")
            run.font.size = Pt(9)
            run.font.color.rgb = resolve_color("primary_light", theme)


def render_ascend(slide, slide_data: dict, theme: dict):
    """Render ascending columns (short→tall) showing growth/phases."""
    y_cursor = render_slide_chrome(slide, slide_data, theme)

    nodes = slide_data.get("nodes", [])
    if not nodes:
        return

    fonts = theme.get("fonts", {})
    body_font = fonts.get("body", {})
    heading_font = fonts.get("heading", {})
    slide_w = theme.get("slide", {}).get("width_in", 13.333)
    bottom = theme.get("margins", {}).get("footer_bottom", 7.0)

    count = len(nodes)
    col_width = 2.0
    gap = 0.25
    positions = distribute_x(count, col_width, gap, slide_width=slide_w)

    max_height = bottom - y_cursor - 1.5
    min_height = max_height * 0.3

    # Palette
    palette = ["accent", "slate_blue", "sage", "gold", "light_blue"]

    for i, node in enumerate(nodes):
        x = positions[i]
        # Height increases linearly
        frac = (i + 1) / count
        col_h = min_height + (max_height - min_height) * frac
        col_top = bottom - 0.5 - col_h
        color = palette[i % len(palette)]

        # Column bar
        add_rect(slide, x, col_top, col_width, col_h, color, theme,
                 corner_radius_emu=45720)

        # Title inside top of bar
        tb = add_textbox(slide, x + 0.1, col_top + 0.15, col_width - 0.2, 0.35)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = node.get("title", "")
        run.font.name = body_font.get("name", "Lato")
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = resolve_color("text_on_dark", theme)

        # Stat label below bar
        stat_text = node.get("stat", "")
        if stat_text:
            tb = add_textbox(slide, x, bottom - 0.35, col_width, 0.3)
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = stat_text
            run.font.name = body_font.get("name", "Lato")
            run.font.size = Pt(9)
            run.font.color.rgb = resolve_color("primary_light", theme)

        # Description inside bar
        desc = node.get("desc", "")
        if desc and col_h > 1.0:
            tb = add_textbox(slide, x + 0.1, col_top + 0.55, col_width - 0.2, col_h - 0.8)
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = desc
            run.font.name = body_font.get("name", "Lato")
            run.font.size = Pt(9)
            run.font.color.rgb = resolve_color("text_on_dark", theme)
