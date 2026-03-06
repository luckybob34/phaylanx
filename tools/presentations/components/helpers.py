"""
Shared helpers for PPTX component renderers.

Provides colour resolution, font application, text run rendering,
shape creation shortcuts, and geometry helpers.
"""

from __future__ import annotations

import math
from typing import Any

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ---------------------------------------------------------------------------
# Colour helpers
# ---------------------------------------------------------------------------

def resolve_color(token: str, theme: dict) -> RGBColor:
    """Resolve a colour token or hex string to an RGBColor.

    Accepts:
      - A hex string like "E55F4C" or "#E55F4C"
      - A token name like "accent" or "primary"
    """
    token = token.strip().lstrip("#")

    # Check if it's already a 6-char hex
    if len(token) == 6:
        try:
            int(token, 16)
            return RGBColor.from_string(token)
        except ValueError:
            pass

    # Look up in theme colors
    colors = theme.get("colors", {})
    palette = theme.get("palette", {})

    hex_val = colors.get(token) or palette.get(token)
    if hex_val:
        return RGBColor.from_string(hex_val.lstrip("#"))

    # Fallback to charcoal
    return RGBColor.from_string("3A3A3A")


def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert a hex string (with or without #) to RGBColor."""
    return RGBColor.from_string(hex_str.strip().lstrip("#"))


# ---------------------------------------------------------------------------
# Font helpers
# ---------------------------------------------------------------------------

def apply_font(run, font_cfg: dict, theme: dict, size_key: str = "size", color_key: str | None = None):
    """Apply font settings from a font config dict to a run's font.

    Args:
        run: pptx Run object
        font_cfg: dict with 'name', 'bold', 'size' etc.
        theme: full theme config for colour resolution
        size_key: which size key to use from font_cfg (e.g. 'size', 'size_large')
        color_key: colour token to apply (overrides font_cfg['color'])
    """
    font = run.font
    font.name = font_cfg.get("name", "Lato")

    if size_key in font_cfg:
        font.size = Pt(font_cfg[size_key])
    elif "size" in font_cfg:
        font.size = Pt(font_cfg["size"])

    if font_cfg.get("bold", False):
        font.bold = True

    color_token = color_key or font_cfg.get("color")
    if color_token:
        font.color.rgb = resolve_color(color_token, theme)


# ---------------------------------------------------------------------------
# Text frame helpers
# ---------------------------------------------------------------------------

def set_text(text_frame, text: str, font_cfg: dict, theme: dict,
             size_key: str = "size", color_key: str | None = None,
             alignment: PP_ALIGN = PP_ALIGN.LEFT, bold: bool | None = None):
    """Set text in a text frame with consistent font styling."""
    text_frame.clear()
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    apply_font(run, font_cfg, theme, size_key, color_key)
    if bold is not None:
        run.font.bold = bold


def add_paragraph(text_frame, text: str, font_cfg: dict, theme: dict,
                  size_key: str = "size", color_key: str | None = None,
                  alignment: PP_ALIGN = PP_ALIGN.LEFT, bold: bool | None = None,
                  level: int = 0, space_before: Pt | None = None,
                  space_after: Pt | None = None):
    """Add a new paragraph to a text frame."""
    p = text_frame.add_paragraph()
    p.alignment = alignment
    p.level = level
    if space_before is not None:
        p.space_before = space_before
    if space_after is not None:
        p.space_after = space_after
    run = p.add_run()
    run.text = text
    apply_font(run, font_cfg, theme, size_key, color_key)
    if bold is not None:
        run.font.bold = bold
    return p


def render_runs(paragraph, runs: list[dict], font_cfg: dict, theme: dict,
                size_key: str = "size", color_key: str | None = None):
    """Render a list of text runs (from parse_outline) into a paragraph."""
    for run_data in runs:
        run = paragraph.add_run()
        run.text = run_data.get("text", "")
        apply_font(run, font_cfg, theme, size_key, color_key)
        if run_data.get("bold"):
            run.font.bold = True
        if run_data.get("italic"):
            run.font.italic = True
        if run_data.get("strike"):
            run.font.strikethrough = True
        if run_data.get("code"):
            mono = theme.get("fonts", {}).get("mono", {})
            run.font.name = mono.get("name", "Consolas")
            if "size" in mono:
                run.font.size = Pt(mono["size"])


def render_body(text_frame, body_paragraphs: list[dict], theme: dict):
    """Render parsed body content (from parse_outline) into a text frame.

    body_paragraphs: list of {"level": int, "runs": [...]}
    """
    text_frame.clear()
    text_frame.word_wrap = True
    body_font = theme.get("fonts", {}).get("body", {})

    first = True
    for para_data in body_paragraphs:
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()

        level = para_data.get("level", 0)
        p.level = level
        if level > 0:
            p.space_before = Pt(2)
            p.space_after = Pt(2)
        else:
            p.space_before = Pt(4)
            p.space_after = Pt(4)

        render_runs(p, para_data.get("runs", []), body_font, theme)


# ---------------------------------------------------------------------------
# Shape creation helpers
# ---------------------------------------------------------------------------

def add_textbox(slide, left_in: float, top_in: float, width_in: float, height_in: float) -> Any:
    """Add a textbox shape and return it."""
    return slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in)
    )


def add_rect(slide, left_in: float, top_in: float, width_in: float, height_in: float,
             fill_token: str, theme: dict, corner_radius_emu: int = 0) -> Any:
    """Add a rounded rectangle shape with solid fill."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius_emu > 0 else MSO_SHAPE.RECTANGLE,
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = resolve_color(fill_token, theme)
    shape.line.fill.background()  # no border

    if corner_radius_emu > 0:
        # Set corner radius via XML adjustment
        try:
            shape.adjustments[0] = corner_radius_emu / shape.width
        except (IndexError, ZeroDivisionError):
            pass

    return shape


def add_oval(slide, left_in: float, top_in: float, width_in: float, height_in: float,
             fill_token: str, theme: dict) -> Any:
    """Add an oval/circle shape with solid fill."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = resolve_color(fill_token, theme)
    shape.line.fill.background()
    return shape


def add_line(slide, start_x_in: float, start_y_in: float,
             end_x_in: float, end_y_in: float,
             color_token: str, theme: dict, width_pt: float = 1.5) -> Any:
    """Add a connector line."""
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(start_x_in), Inches(start_y_in),
        Inches(end_x_in), Inches(end_y_in),
    )
    connector.line.color.rgb = resolve_color(color_token, theme)
    connector.line.width = Pt(width_pt)
    return connector


def add_arrow_right(slide, left_in: float, top_in: float, width_in: float, height_in: float,
                    color_token: str, theme: dict) -> Any:
    """Add a right-pointing arrow shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = resolve_color(color_token, theme)
    shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# Geometry helpers
# ---------------------------------------------------------------------------

def center_x(width_in: float, slide_width: float = 13.333) -> float:
    """Calculate left position to centre a shape horizontally."""
    return (slide_width - width_in) / 2


def distribute_x(count: int, item_width_in: float, gap_in: float,
                  margin_left: float = 0.66, slide_width: float = 13.333) -> list[float]:
    """Calculate evenly distributed X positions for N items."""
    total = count * item_width_in + (count - 1) * gap_in
    start = (slide_width - total) / 2
    return [start + i * (item_width_in + gap_in) for i in range(count)]


def point_on_circle(cx: float, cy: float, radius: float, angle_deg: float) -> tuple[float, float]:
    """Calculate a point on a circle given centre, radius, and angle in degrees."""
    rad = math.radians(angle_deg - 90)  # Start from top (12 o'clock)
    return (cx + radius * math.cos(rad), cy + radius * math.sin(rad))
