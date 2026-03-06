"""
Layout renderers — content slides with body text and two-column layouts.
"""

from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from .helpers import (
    add_textbox, resolve_color, render_body, render_runs,
    apply_font, set_text,
)
from .typography import render_slide_chrome


def render_content(slide, slide_data: dict, theme: dict):
    """Render a content slide: title + subtitle + body text."""
    y_cursor = render_slide_chrome(slide, slide_data, theme)

    margins = theme.get("margins", {})
    left = margins.get("left", 0.66)
    slide_w = theme.get("slide", {}).get("width_in", 13.333)
    right = margins.get("right", 0.66)
    width = slide_w - left - right
    bottom = margins.get("footer_bottom", 7.0)
    body_height = bottom - y_cursor - 0.3

    if "body" in slide_data:
        tb = add_textbox(slide, left, y_cursor, width, body_height)
        tf = tb.text_frame
        render_body(tf, slide_data["body"], theme)


def render_two_col(slide, slide_data: dict, theme: dict):
    """Render a two-column layout slide."""
    y_cursor = render_slide_chrome(slide, slide_data, theme)

    margins = theme.get("margins", {})
    left = margins.get("left", 0.66)
    slide_w = theme.get("slide", {}).get("width_in", 13.333)
    right = margins.get("right", 0.66)
    total_width = slide_w - left - right
    gap = 0.4
    col_width = (total_width - gap) / 2
    bottom = margins.get("footer_bottom", 7.0)
    col_height = bottom - y_cursor - 0.3

    # Left column
    if "left" in slide_data:
        tb = add_textbox(slide, left, y_cursor, col_width, col_height)
        tf = tb.text_frame
        render_body(tf, slide_data["left"], theme)

    # Right column
    if "right" in slide_data:
        tb = add_textbox(slide, left + col_width + gap, y_cursor, col_width, col_height)
        tf = tb.text_frame
        render_body(tf, slide_data["right"], theme)
