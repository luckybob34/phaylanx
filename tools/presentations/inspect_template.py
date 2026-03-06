#!/usr/bin/env python3
"""inspect_template.py — Enumerate every slide master, layout, placeholder, and
theme colour inside a .pptx template.  Outputs structured JSON to stdout (or a
file with --out).

Usage:
    python inspect_template.py <path-to-template.pptx> [--out report.json]

The report is the primary input for building the PPTX theme `config.yaml`.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN


def _emu_to_inches(emu: int | None) -> float | None:
    return round(emu / 914400, 3) if emu is not None else None


def _placeholder_info(ph) -> dict:
    """Extract useful details from a placeholder shape."""
    return {
        "idx": ph.placeholder_format.idx,
        "type": str(ph.placeholder_format.type),
        "name": ph.name,
        "left_in": _emu_to_inches(ph.left),
        "top_in": _emu_to_inches(ph.top),
        "width_in": _emu_to_inches(ph.width),
        "height_in": _emu_to_inches(ph.height),
    }


def _shape_info(shape) -> dict:
    """Extract useful details from any shape."""
    info = {
        "shape_id": shape.shape_id,
        "name": shape.name,
        "shape_type": str(shape.shape_type) if shape.shape_type else None,
        "left_in": _emu_to_inches(shape.left),
        "top_in": _emu_to_inches(shape.top),
        "width_in": _emu_to_inches(shape.width),
        "height_in": _emu_to_inches(shape.height),
    }
    if shape.has_text_frame:
        info["text_preview"] = shape.text_frame.text[:120] if shape.text_frame.text else ""
    try:
        pf = shape.placeholder_format
        if pf is not None:
            info["placeholder"] = _placeholder_info(shape)
    except (ValueError, AttributeError):
        pass
    return info


def _extract_theme_colors(prs: Presentation) -> dict:
    """Pull colour scheme from the theme XML."""
    colors = {}
    try:
        theme_el = prs.slide_masters[0].element
        # Navigate to the theme element via the relationship
        from lxml import etree
        theme_part = prs.slide_masters[0].part.slide_master.part
        # Alternative: reach into the package parts
        for rel in prs.slide_masters[0].part.rels.values():
            if "theme" in rel.reltype:
                theme_xml = rel.target_part.blob
                root = etree.fromstring(theme_xml)
                ns = {
                    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                }
                # Color scheme
                clr_scheme = root.find(".//a:themeElements/a:clrScheme", ns)
                if clr_scheme is not None:
                    colors["scheme_name"] = clr_scheme.get("name", "")
                    for child in clr_scheme:
                        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                        for val in child:
                            val_tag = val.tag.split("}")[-1] if "}" in val.tag else val.tag
                            if val_tag == "srgbClr":
                                colors[tag] = f"#{val.get('val', '')}"
                            elif val_tag == "sysClr":
                                colors[tag] = val.get("lastClr", val.get("val", ""))
                # Font scheme
                font_scheme = root.find(".//a:themeElements/a:fontScheme", ns)
                if font_scheme is not None:
                    colors["font_scheme_name"] = font_scheme.get("name", "")
                    major = font_scheme.find("a:majorFont/a:latin", ns)
                    minor = font_scheme.find("a:minorFont/a:latin", ns)
                    if major is not None:
                        colors["major_font"] = major.get("typeface", "")
                    if minor is not None:
                        colors["minor_font"] = minor.get("typeface", "")
                break
    except Exception as e:
        colors["_error"] = str(e)
    return colors


def inspect(pptx_path: str) -> dict:
    """Return a full inspection report of the PPTX template."""
    prs = Presentation(pptx_path)

    report: dict = {
        "file": str(pptx_path),
        "slide_width_in": _emu_to_inches(prs.slide_width),
        "slide_height_in": _emu_to_inches(prs.slide_height),
        "slide_count": len(prs.slides),
        "slide_masters": [],
        "theme_colors": {},
    }

    # ----- Slide masters & layouts -----
    for m_idx, master in enumerate(prs.slide_masters):
        master_info = {
            "master_index": m_idx,
            "name": master.name if hasattr(master, "name") else f"Master {m_idx}",
            "shapes": [_shape_info(s) for s in master.shapes],
            "layouts": [],
        }
        for l_idx, layout in enumerate(master.slide_layouts):
            layout_info = {
                "layout_index": l_idx,
                "name": layout.name,
                "placeholders": [],
                "shapes": [],
            }
            for ph in layout.placeholders:
                layout_info["placeholders"].append(_placeholder_info(ph))
            for shape in layout.shapes:
                try:
                    _ = shape.placeholder_format
                    # it's a placeholder — already captured above
                except (ValueError, AttributeError):
                    layout_info["shapes"].append(_shape_info(shape))
            master_info["layouts"].append(layout_info)
        report["slide_masters"].append(master_info)

    # ----- Existing slides -----
    slides_info = []
    for s_idx, slide in enumerate(prs.slides):
        slide_info = {
            "slide_index": s_idx,
            "layout_name": slide.slide_layout.name,
            "shapes": [_shape_info(s) for s in slide.shapes],
        }
        slides_info.append(slide_info)
    report["existing_slides"] = slides_info

    # ----- Theme colours -----
    report["theme_colors"] = _extract_theme_colors(prs)

    return report


def main():
    parser = argparse.ArgumentParser(description="Inspect a PPTX template")
    parser.add_argument("pptx", help="Path to the .pptx file")
    parser.add_argument("--out", "-o", help="Write JSON report to this file (default: stdout)")
    args = parser.parse_args()

    if not Path(args.pptx).exists():
        print(f"Error: {args.pptx} not found", file=sys.stderr)
        sys.exit(1)

    report = inspect(args.pptx)
    output = json.dumps(report, indent=2, ensure_ascii=False)

    if args.out:
        Path(args.out).write_text(output, encoding="utf-8")
        print(f"Report written to {args.out}")
    else:
        print(output)


if __name__ == "__main__":
    main()
