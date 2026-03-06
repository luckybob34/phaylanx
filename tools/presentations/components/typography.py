"""
Typography renderer — places title, subtitle, eyebrow, and speaker notes.

Called by render_pptx.py for every slide before the component renderer runs.
"""

from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .helpers import (
    add_textbox, resolve_color, apply_font, set_text, render_runs,
)


def render_slide_chrome(slide, slide_data: dict, theme: dict, *, is_dark: bool = False):
    """Render the standard slide chrome: eyebrow, title, subtitle.

    Args:
        slide: pptx Slide object
        slide_data: parsed slide dict from outline
        theme: full theme config
        is_dark: if True, use light text colours
    """
    margins = theme.get("margins", {})
    fonts = theme.get("fonts", {})
    heading_font = fonts.get("heading", {})
    body_font = fonts.get("body", {})

    left = margins.get("left", 0.66)
    right_edge = theme.get("slide", {}).get("width_in", 13.333) - margins.get("right", 0.66)
    width = right_edge - left

    color_key = "text_on_dark" if is_dark else None
    y_cursor = margins.get("top", 0.67)

    # --- Eyebrow ---
    if "eyebrow" in slide_data:
        tb = add_textbox(slide, left, y_cursor, width, 0.25)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = slide_data["eyebrow"].upper()
        run.font.name = body_font.get("name", "Lato")
        run.font.size = Pt(body_font.get("size_label", 10))
        run.font.bold = True
        run.font.color.rgb = resolve_color(
            color_key or "accent", theme
        )
        # Tracking (letter spacing) via XML
        run.font._element.attrib[
            "{http://schemas.openxmlformats.org/drawingml/2006/main}spc"
        ] = "150"  # 1.5pt spacing
        y_cursor += 0.35

    # --- Title ---
    if "title" in slide_data:
        layout = slide_data.get("layout", "content")
        if layout == "title":
            size_key = "size_hero"
        elif layout == "section":
            size_key = "size_section"
        else:
            size_key = "size_slide"

        title_height = 0.6 if size_key == "size_slide" else 0.9
        tb = add_textbox(slide, left, y_cursor, width, title_height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = slide_data["title"]
        apply_font(run, heading_font, theme, size_key, color_key)
        y_cursor += title_height + 0.1

    # --- Subtitle ---
    if "subtitle" in slide_data:
        tb = add_textbox(slide, left, y_cursor, width, 0.5)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = slide_data["subtitle"]
        apply_font(run, heading_font, theme, "size_subtitle", color_key or "primary_light")
        y_cursor += 0.6

    # --- Speaker notes ---
    if "notes" in slide_data:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = slide_data["notes"]

    return y_cursor


def render_title_slide(slide, slide_data: dict, theme: dict):
    """Render a cover/title slide with centred hero title."""
    fonts = theme.get("fonts", {})
    heading = fonts.get("heading", {})
    slide_w = theme.get("slide", {}).get("width_in", 13.333)
    slide_h = theme.get("slide", {}).get("height_in", 7.5)

    # Hero title — centred vertically and horizontally
    title_width = slide_w * 0.5
    title_left = (slide_w - title_width) / 2
    title_top = slide_h * 0.28

    if "title" in slide_data:
        tb = add_textbox(slide, title_left - 1.5, title_top, title_width + 3, 1.8)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = slide_data["title"]
        apply_font(run, heading, theme, "size_hero", "text_on_dark")

    # Subtitle below
    if "subtitle" in slide_data:
        tb = add_textbox(slide, title_left - 1.5, title_top + 1.8, title_width + 3, 0.5)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = slide_data["subtitle"]
        apply_font(run, heading, theme, "size_subtitle", "text_on_dark")

    # Speaker notes
    if "notes" in slide_data:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = slide_data["notes"]


def render_section_slide(slide, slide_data: dict, theme: dict):
    """Render a section break slide with centred title."""
    fonts = theme.get("fonts", {})
    heading = fonts.get("heading", {})
    slide_w = theme.get("slide", {}).get("width_in", 13.333)
    slide_h = theme.get("slide", {}).get("height_in", 7.5)

    title_width = 6.0
    title_left = (slide_w - title_width) / 2
    title_top = slide_h * 0.42

    if "title" in slide_data:
        tb = add_textbox(slide, title_left, title_top, title_width, 0.8)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = slide_data["title"]
        apply_font(run, heading, theme, "size_section", "text_on_dark")

    if "notes" in slide_data:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = slide_data["notes"]


def render_stat_slide(slide, slide_data: dict, theme: dict):
    """Render an impressive stat slide with a big number."""
    fonts = theme.get("fonts", {})
    stat_font = fonts.get("stat", {})
    heading = fonts.get("heading", {})
    body = fonts.get("body", {})
    slide_w = theme.get("slide", {}).get("width_in", 13.333)
    slide_h = theme.get("slide", {}).get("height_in", 7.5)

    # Big number centred
    if "value" in slide_data:
        val_width = 8.0
        val_left = (slide_w - val_width) / 2
        tb = add_textbox(slide, val_left, 2.2, val_width, 1.8)
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = slide_data["value"]
        apply_font(run, stat_font, theme, "size_large")

    # Title below
    if "title" in slide_data:
        tb = add_textbox(slide, (slide_w - 8) / 2, 4.0, 8, 0.7)
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = slide_data["title"]
        apply_font(run, heading, theme, "size_slide", "text_on_dark")

    # Label / subtitle
    label = slide_data.get("label") or slide_data.get("subtitle")
    if label:
        tb = add_textbox(slide, (slide_w - 8) / 2, 4.8, 8, 0.5)
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        apply_font(run, body, theme, "size", "text_on_dark")

    if "notes" in slide_data:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = slide_data["notes"]


def render_quote_slide(slide, slide_data: dict, theme: dict):
    """Render a quote slide."""
    fonts = theme.get("fonts", {})
    heading = fonts.get("heading", {})
    body = fonts.get("body", {})
    margins = theme.get("margins", {})
    left = margins.get("left", 0.66)
    slide_w = theme.get("slide", {}).get("width_in", 13.333)
    width = slide_w - 2 * left

    # Quote mark
    tb = add_textbox(slide, left, 1.8, 1.0, 1.0)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "\u201C"  # Opening double quote
    run.font.size = Pt(72)
    run.font.color.rgb = resolve_color("accent", theme)
    run.font.name = heading.get("name", "Georgia")

    # Quote text
    if "quote" in slide_data:
        tb = add_textbox(slide, left + 0.5, 2.8, width - 1.0, 2.5)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = slide_data["quote"]
        run.font.name = heading.get("name", "Georgia")
        run.font.size = Pt(18)
        run.font.italic = True
        run.font.color.rgb = resolve_color("primary", theme)

    # Attribution
    y = 5.5
    if "author" in slide_data:
        tb = add_textbox(slide, left + 0.5, y, width - 1.0, 0.4)
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"— {slide_data['author']}"
        run.font.name = body.get("name", "Lato")
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = resolve_color("primary", theme)
        y += 0.35

    if "role" in slide_data:
        tb = add_textbox(slide, left + 0.5, y, width - 1.0, 0.3)
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = slide_data["role"]
        run.font.name = body.get("name", "Lato")
        run.font.size = Pt(11)
        run.font.color.rgb = resolve_color("primary_light", theme)

    if "notes" in slide_data:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = slide_data["notes"]
